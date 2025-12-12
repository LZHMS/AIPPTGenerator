"""
PPT Builder Module
使用 python-pptx 生成实际的 PPT 文件
支持多种主题风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from typing import Dict, Any, List
import re


def hex_to_rgb(hex_color: str) -> tuple:
    """将十六进制颜色转换为RGB元组"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_rgb_color(hex_color: str) -> RGBColor:
    """从十六进制颜色创建RGBColor对象"""
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


class PPTBuilder:
    """PPT 构建器"""
    
    def __init__(self, ppt_data: Dict[str, Any]):
        """
        初始化 PPT 构建器
        
        Args:
            ppt_data: PPT 数据字典，包含主题、配色、幻灯片内容等
        """
        self.ppt_data = ppt_data
        self.prs = Presentation()
        
        # 设置幻灯片尺寸（16:9）
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # 获取配色方案
        self.colors = ppt_data.get("color_scheme", {})
        self.theme = ppt_data.get("theme_style", {})
        
        # 默认配色
        self.default_colors = {
            "primary_color": "#2E86AB",
            "secondary_color": "#A23B72",
            "accent_color": "#F18F01",
            "background_color": "#FFFFFF",
            "text_color": "#333333",
            "title_color": "#1A1A2E",
            "gradient_start": "#667eea",
            "gradient_end": "#764ba2"
        }
        
        # 合并配色
        for key, value in self.default_colors.items():
            if key not in self.colors or not self.colors[key]:
                self.colors[key] = value
    
    def get_color(self, color_key: str) -> str:
        """获取颜色值"""
        return self.colors.get(color_key, self.default_colors.get(color_key, "#333333"))
    
    def add_background(self, slide, color: str = None):
        """为幻灯片添加背景色"""
        if color is None:
            color = self.get_color("background_color")
        
        # 添加背景形状
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = create_rgb_color(color)
        background.line.fill.background()
        
        # 将背景移到最底层
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    def add_title_slide(self, slide_data: Dict):
        """添加标题页"""
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加背景
        self.add_background(slide)
        
        # 添加装饰性元素
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(3.2),
            Inches(13.333), Inches(0.1)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
        accent_shape.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        if slide_data.get("subtitle") or slide_data.get("content"):
            subtitle_text = slide_data.get("subtitle", "")
            if not subtitle_text and slide_data.get("content"):
                subtitle_text = slide_data["content"][0] if slide_data["content"] else ""
            
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.8),
                Inches(12.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = subtitle_text
            subtitle_p.font.size = Pt(28)
            subtitle_p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
            subtitle_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(self, slide_data: Dict):
        """添加内容页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加背景
        self.add_background(slide)
        
        # 添加顶部装饰条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.8)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
        top_bar.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 添加内容
        content_items = slide_data.get("content", [])
        if content_items:
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.3),
                Inches(11.733), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for idx, item in enumerate(content_items):
                if idx == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.font.size = Pt(24)
                p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
                p.space_after = Pt(18)
                p.level = 0
        
        return slide
    
    def add_bullet_slide(self, slide_data: Dict):
        """添加要点列表页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加背景
        self.add_background(slide)
        
        # 添加左侧装饰条
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = create_rgb_color(self.get_color("accent_color"))
        left_bar.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5),
            Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        
        # 添加分隔线
        sep_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.5),
            Inches(3), Inches(0.05)
        )
        sep_line.fill.solid()
        sep_line.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
        sep_line.line.fill.background()
        
        # 添加要点
        content_items = slide_data.get("content", [])
        if content_items:
            start_y = 2.0
            for idx, item in enumerate(content_items):
                # 添加要点编号圆形
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(0.8), Inches(start_y + idx * 1.2),
                    Inches(0.4), Inches(0.4)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
                circle.line.fill.background()
                
                # 添加编号文字
                num_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(start_y + idx * 1.2),
                    Inches(0.4), Inches(0.4)
                )
                num_frame = num_box.text_frame
                num_frame.paragraphs[0].text = str(idx + 1)
                num_frame.paragraphs[0].font.size = Pt(18)
                num_frame.paragraphs[0].font.bold = True
                num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # 添加要点文字
                item_box = slide.shapes.add_textbox(
                    Inches(1.5), Inches(start_y + idx * 1.2 - 0.1),
                    Inches(11), Inches(1)
                )
                item_frame = item_box.text_frame
                item_frame.word_wrap = True
                item_p = item_frame.paragraphs[0]
                item_p.text = item
                item_p.font.size = Pt(24)
                item_p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
        
        return slide
    
    def add_two_column_slide(self, slide_data: Dict):
        """添加两栏布局页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加背景
        self.add_background(slide)
        
        # 添加标题背景
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
        title_bg.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER
        
        # 分割内容到两栏
        content_items = slide_data.get("content", [])
        mid = len(content_items) // 2
        left_items = content_items[:max(mid, 1)]
        right_items = content_items[max(mid, 1):]
        
        # 左栏
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(5.5)
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = create_rgb_color("#F8F9FA")
        left_box.line.color.rgb = create_rgb_color(self.get_color("primary_color"))
        
        left_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8),
            Inches(5.4), Inches(5)
        )
        left_frame = left_text.text_frame
        left_frame.word_wrap = True
        for idx, item in enumerate(left_items):
            if idx == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
            p.space_after = Pt(12)
        
        # 右栏
        if right_items:
            right_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.833), Inches(1.5),
                Inches(6), Inches(5.5)
            )
            right_box.fill.solid()
            right_box.fill.fore_color.rgb = create_rgb_color("#F8F9FA")
            right_box.line.color.rgb = create_rgb_color(self.get_color("secondary_color"))
            
            right_text = slide.shapes.add_textbox(
                Inches(7.133), Inches(1.8),
                Inches(5.4), Inches(5)
            )
            right_frame = right_text.text_frame
            right_frame.word_wrap = True
            for idx, item in enumerate(right_items):
                if idx == 0:
                    p = right_frame.paragraphs[0]
                else:
                    p = right_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(20)
                p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
                p.space_after = Pt(12)
        
        return slide
    
    def add_summary_slide(self, slide_data: Dict):
        """添加总结页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加渐变背景效果（使用纯色近似）
        self.add_background(slide, self.get_color("primary_color"))
        
        # 添加装饰圆形
        circle1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-2), Inches(-2),
            Inches(6), Inches(6)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = create_rgb_color(self.get_color("secondary_color"))
        circle1.fill.fore_color.brightness = 0.3
        circle1.line.fill.background()
        
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(10), Inches(4),
            Inches(5), Inches(5)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = create_rgb_color(self.get_color("accent_color"))
        circle2.fill.fore_color.brightness = 0.3
        circle2.line.fill.background()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "总结")
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER
        
        # 添加内容
        content_items = slide_data.get("content", [])
        if content_items:
            content_text = " | ".join(content_items)
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5),
                Inches(12.333), Inches(2)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_p = content_frame.paragraphs[0]
            content_p.text = content_text
            content_p.font.size = Pt(24)
            content_p.font.color.rgb = RGBColor(255, 255, 255)
            content_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_three_column_slide(self, slide_data: Dict):
        """添加三栏布局页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 分割内容到三栏
        content_items = slide_data.get("content", [])
        cols = [[], [], []]
        for i, item in enumerate(content_items):
            cols[i % 3].append(item)
        
        col_width = 4.0
        col_start_x = [0.5, 4.7, 8.9]
        
        for col_idx, col_items in enumerate(cols):
            if not col_items:
                continue
            
            # 列标题背景
            col_header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(col_start_x[col_idx]), Inches(1.3),
                Inches(col_width), Inches(0.6)
            )
            col_header.fill.solid()
            colors = [self.get_color("primary_color"), self.get_color("secondary_color"), self.get_color("accent_color")]
            col_header.fill.fore_color.rgb = create_rgb_color(colors[col_idx])
            col_header.line.fill.background()
            
            # 列内容区域
            col_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(col_start_x[col_idx]), Inches(2.0),
                Inches(col_width), Inches(5.0)
            )
            col_box.fill.solid()
            col_box.fill.fore_color.rgb = create_rgb_color("#F8F9FA")
            col_box.line.color.rgb = create_rgb_color(colors[col_idx])
            
            # 列内容文字
            col_text = slide.shapes.add_textbox(
                Inches(col_start_x[col_idx] + 0.2), Inches(2.2),
                Inches(col_width - 0.4), Inches(4.6)
            )
            col_frame = col_text.text_frame
            col_frame.word_wrap = True
            
            for idx, item in enumerate(col_items):
                if idx == 0:
                    p = col_frame.paragraphs[0]
                else:
                    p = col_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
                p.space_after = Pt(10)
        
        return slide
    
    def add_quote_slide(self, slide_data: Dict):
        """添加引用/名言页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide, self.get_color("primary_color"))
        
        # 添加引号装饰
        quote_mark = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(2), Inches(2)
        )
        quote_frame = quote_mark.text_frame
        quote_p = quote_frame.paragraphs[0]
        quote_p.text = '"'
        quote_p.font.size = Pt(150)
        quote_p.font.color.rgb = RGBColor(255, 255, 255)
        quote_p.font.bold = True
        
        # 引用内容
        content_items = slide_data.get("content", [])
        quote_text = content_items[0] if content_items else slide_data.get("title", "")
        
        content_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5),
            Inches(9.333), Inches(3)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_p = content_frame.paragraphs[0]
        content_p.text = quote_text
        content_p.font.size = Pt(32)
        content_p.font.italic = True
        content_p.font.color.rgb = RGBColor(255, 255, 255)
        content_p.alignment = PP_ALIGN.CENTER
        
        # 来源
        if len(content_items) > 1:
            source_box = slide.shapes.add_textbox(
                Inches(2), Inches(5.5),
                Inches(9.333), Inches(1)
            )
            source_frame = source_box.text_frame
            source_p = source_frame.paragraphs[0]
            source_p.text = f"— {content_items[1]}"
            source_p.font.size = Pt(20)
            source_p.font.color.rgb = RGBColor(200, 200, 200)
            source_p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def add_statistics_slide(self, slide_data: Dict):
        """添加数据统计页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "数据统计")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 数据卡片
        content_items = slide_data.get("content", [])
        num_items = min(len(content_items), 4)
        card_width = 2.8
        total_width = num_items * card_width + (num_items - 1) * 0.3
        start_x = (13.333 - total_width) / 2
        
        colors = [self.get_color("primary_color"), self.get_color("secondary_color"), 
                  self.get_color("accent_color"), "#28a745"]
        
        for idx, item in enumerate(content_items[:4]):
            card_x = start_x + idx * (card_width + 0.3)
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(card_x), Inches(2),
                Inches(card_width), Inches(4.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = create_rgb_color(colors[idx % len(colors)])
            card.line.fill.background()
            
            # 卡片内容
            card_text = slide.shapes.add_textbox(
                Inches(card_x + 0.2), Inches(2.5),
                Inches(card_width - 0.4), Inches(3.5)
            )
            card_frame = card_text.text_frame
            card_frame.word_wrap = True
            
            p = card_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_timeline_slide(self, slide_data: Dict):
        """添加时间线布局页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "时间线")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 时间线主轴
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(3.75),
            Inches(11.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
        line.line.fill.background()
        
        # 时间点
        content_items = slide_data.get("content", [])
        num_items = min(len(content_items), 5)
        if num_items > 0:
            spacing = 11.333 / (num_items + 1)
            
            for idx, item in enumerate(content_items[:5]):
                x_pos = 1 + spacing * (idx + 1)
                
                # 圆点
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x_pos - 0.15), Inches(3.6),
                    Inches(0.3), Inches(0.3)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = create_rgb_color(self.get_color("accent_color"))
                dot.line.fill.background()
                
                # 上方或下方显示文字（交替）
                y_pos = 1.8 if idx % 2 == 0 else 4.3
                
                text_box = slide.shapes.add_textbox(
                    Inches(x_pos - 1), Inches(y_pos),
                    Inches(2), Inches(1.5)
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_comparison_slide(self, slide_data: Dict):
        """添加对比布局页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "对比分析")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # VS 标记
        vs_box = slide.shapes.add_textbox(
            Inches(6.166), Inches(3.2),
            Inches(1), Inches(1)
        )
        vs_frame = vs_box.text_frame
        vs_p = vs_frame.paragraphs[0]
        vs_p.text = "VS"
        vs_p.font.size = Pt(36)
        vs_p.font.bold = True
        vs_p.font.color.rgb = create_rgb_color(self.get_color("accent_color"))
        vs_p.alignment = PP_ALIGN.CENTER
        
        # 分割内容
        content_items = slide_data.get("content", [])
        mid = len(content_items) // 2
        left_items = content_items[:max(mid, 1)]
        right_items = content_items[max(mid, 1):]
        
        # 左侧
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.3),
            Inches(5.5), Inches(5.7)
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = create_rgb_color(self.get_color("primary_color"))
        left_box.line.fill.background()
        
        left_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6),
            Inches(4.9), Inches(5.1)
        )
        left_frame = left_text.text_frame
        left_frame.word_wrap = True
        for idx, item in enumerate(left_items):
            if idx == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = f"✓ {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
        
        # 右侧
        right_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.333), Inches(1.3),
            Inches(5.5), Inches(5.7)
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = create_rgb_color(self.get_color("secondary_color"))
        right_box.line.fill.background()
        
        right_text = slide.shapes.add_textbox(
            Inches(7.633), Inches(1.6),
            Inches(4.9), Inches(5.1)
        )
        right_frame = right_text.text_frame
        right_frame.word_wrap = True
        for idx, item in enumerate(right_items):
            if idx == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = f"✓ {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
        
        return slide
    
    def add_icons_grid_slide(self, slide_data: Dict):
        """添加图标网格页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 网格布局
        content_items = slide_data.get("content", [])
        num_items = min(len(content_items), 6)
        
        # 2行3列布局
        cols = 3
        rows = 2
        cell_width = 3.8
        cell_height = 2.5
        start_x = (13.333 - cols * cell_width - 0.5) / 2
        start_y = 1.5
        
        icons = ["★", "◆", "●", "▲", "■", "◉"]
        colors = [self.get_color("primary_color"), self.get_color("secondary_color"), 
                  self.get_color("accent_color")] * 2
        
        for idx, item in enumerate(content_items[:6]):
            row = idx // cols
            col = idx % cols
            x = start_x + col * (cell_width + 0.25)
            y = start_y + row * (cell_height + 0.3)
            
            # 图标圆形背景
            icon_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + cell_width/2 - 0.4), Inches(y),
                Inches(0.8), Inches(0.8)
            )
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = create_rgb_color(colors[idx])
            icon_bg.line.fill.background()
            
            # 图标
            icon_box = slide.shapes.add_textbox(
                Inches(x + cell_width/2 - 0.4), Inches(y + 0.1),
                Inches(0.8), Inches(0.6)
            )
            icon_frame = icon_box.text_frame
            icon_p = icon_frame.paragraphs[0]
            icon_p.text = icons[idx % len(icons)]
            icon_p.font.size = Pt(24)
            icon_p.font.color.rgb = RGBColor(255, 255, 255)
            icon_p.alignment = PP_ALIGN.CENTER
            
            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 1),
                Inches(cell_width), Inches(1.3)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_p = text_frame.paragraphs[0]
            text_p.text = item
            text_p.font.size = Pt(18)
            text_p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
            text_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_big_number_slide(self, slide_data: Dict):
        """添加大数字展示页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 大数字
        content_items = slide_data.get("content", [])
        big_text = content_items[0] if content_items else "100%"
        
        number_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(12.333), Inches(3)
        )
        number_frame = number_box.text_frame
        number_p = number_frame.paragraphs[0]
        number_p.text = big_text
        number_p.font.size = Pt(120)
        number_p.font.bold = True
        number_p.font.color.rgb = create_rgb_color(self.get_color("primary_color"))
        number_p.alignment = PP_ALIGN.CENTER
        
        # 描述文字
        if len(content_items) > 1:
            desc_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.5),
                Inches(12.333), Inches(1.5)
            )
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_p = desc_frame.paragraphs[0]
            desc_p.text = " | ".join(content_items[1:])
            desc_p.font.size = Pt(24)
            desc_p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
            desc_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_process_flow_slide(self, slide_data: Dict):
        """添加流程图页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "流程")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        title_p.alignment = PP_ALIGN.CENTER
        
        # 流程步骤
        content_items = slide_data.get("content", [])
        num_items = min(len(content_items), 5)
        
        if num_items > 0:
            step_width = 2.2
            arrow_width = 0.5
            total_width = num_items * step_width + (num_items - 1) * arrow_width
            start_x = (13.333 - total_width) / 2
            
            colors = [self.get_color("primary_color"), self.get_color("secondary_color"),
                      self.get_color("accent_color"), "#28a745", "#6c757d"]
            
            for idx, item in enumerate(content_items[:5]):
                x = start_x + idx * (step_width + arrow_width)
                
                # 步骤圆形
                step_circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x), Inches(2.5),
                    Inches(step_width), Inches(step_width)
                )
                step_circle.fill.solid()
                step_circle.fill.fore_color.rgb = create_rgb_color(colors[idx % len(colors)])
                step_circle.line.fill.background()
                
                # 步骤编号
                num_box = slide.shapes.add_textbox(
                    Inches(x), Inches(2.8),
                    Inches(step_width), Inches(0.8)
                )
                num_frame = num_box.text_frame
                num_p = num_frame.paragraphs[0]
                num_p.text = str(idx + 1)
                num_p.font.size = Pt(36)
                num_p.font.bold = True
                num_p.font.color.rgb = RGBColor(255, 255, 255)
                num_p.alignment = PP_ALIGN.CENTER
                
                # 步骤文字
                text_box = slide.shapes.add_textbox(
                    Inches(x - 0.3), Inches(5),
                    Inches(step_width + 0.6), Inches(1.5)
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_p = text_frame.paragraphs[0]
                text_p.text = item
                text_p.font.size = Pt(16)
                text_p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
                text_p.alignment = PP_ALIGN.CENTER
                
                # 箭头（除了最后一个）
                if idx < num_items - 1:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(x + step_width + 0.1), Inches(3.4),
                        Inches(arrow_width - 0.2), Inches(0.4)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = create_rgb_color(self.get_color("text_color"))
                    arrow.line.fill.background()
        
        return slide
    
    def add_image_left_slide(self, slide_data: Dict):
        """添加左图右文页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        
        # 左侧图片占位区域
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.3),
            Inches(5.5), Inches(5.7)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = create_rgb_color("#E9ECEF")
        img_placeholder.line.color.rgb = create_rgb_color(self.get_color("primary_color"))
        
        # 图片占位文字
        img_text = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.5),
            Inches(3.5), Inches(1)
        )
        img_frame = img_text.text_frame
        img_p = img_frame.paragraphs[0]
        img_p.text = "📷 图片区域"
        img_p.font.size = Pt(24)
        img_p.font.color.rgb = create_rgb_color("#6c757d")
        img_p.alignment = PP_ALIGN.CENTER
        
        # 右侧文字
        content_items = slide_data.get("content", [])
        content_box = slide.shapes.add_textbox(
            Inches(6.5), Inches(1.5),
            Inches(6.333), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for idx, item in enumerate(content_items):
            if idx == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(22)
            p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
            p.space_after = Pt(15)
        
        return slide
    
    def add_image_right_slide(self, slide_data: Dict):
        """添加右图左文页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_background(slide)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.get("title", "")
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = create_rgb_color(self.get_color("title_color"))
        
        # 左侧文字
        content_items = slide_data.get("content", [])
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(6.333), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for idx, item in enumerate(content_items):
            if idx == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(22)
            p.font.color.rgb = create_rgb_color(self.get_color("text_color"))
            p.space_after = Pt(15)
        
        # 右侧图片占位区域
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.333), Inches(1.3),
            Inches(5.5), Inches(5.7)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = create_rgb_color("#E9ECEF")
        img_placeholder.line.color.rgb = create_rgb_color(self.get_color("primary_color"))
        
        # 图片占位文字
        img_text = slide.shapes.add_textbox(
            Inches(8.333), Inches(3.5),
            Inches(3.5), Inches(1)
        )
        img_frame = img_text.text_frame
        img_p = img_frame.paragraphs[0]
        img_p.text = "📷 图片区域"
        img_p.font.size = Pt(24)
        img_p.font.color.rgb = create_rgb_color("#6c757d")
        img_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def build(self) -> Presentation:
        """构建完整的PPT"""
        slides = self.ppt_data.get("slides", [])
        
        for slide_data in slides:
            slide_type = slide_data.get("slide_type", "content")
            
            if slide_type == "title":
                self.add_title_slide(slide_data)
            elif slide_type == "bullet_points":
                self.add_bullet_slide(slide_data)
            elif slide_type == "two_column":
                self.add_two_column_slide(slide_data)
            elif slide_type == "three_column":
                self.add_three_column_slide(slide_data)
            elif slide_type == "quote":
                self.add_quote_slide(slide_data)
            elif slide_type == "statistics":
                self.add_statistics_slide(slide_data)
            elif slide_type == "timeline":
                self.add_timeline_slide(slide_data)
            elif slide_type == "comparison":
                self.add_comparison_slide(slide_data)
            elif slide_type == "icons_grid":
                self.add_icons_grid_slide(slide_data)
            elif slide_type == "big_number":
                self.add_big_number_slide(slide_data)
            elif slide_type == "process_flow":
                self.add_process_flow_slide(slide_data)
            elif slide_type == "image_left":
                self.add_image_left_slide(slide_data)
            elif slide_type == "image_right":
                self.add_image_right_slide(slide_data)
            elif slide_type == "summary":
                self.add_summary_slide(slide_data)
            else:
                self.add_content_slide(slide_data)
        
        return self.prs
    
    def save(self, filepath: str):
        """保存PPT文件"""
        self.build()
        
        # 确保目录存在
        os.makedirs(os.path.dirname(filepath) if os.path.dirname(filepath) else ".", exist_ok=True)
        
        self.prs.save(filepath)
        return filepath


def create_ppt_from_data(ppt_data: Dict[str, Any], output_path: str) -> str:
    """
    从PPT数据创建PPT文件
    
    Args:
        ppt_data: PPT数据字典
        output_path: 输出文件路径
        
    Returns:
        保存的文件路径
    """
    builder = PPTBuilder(ppt_data)
    return builder.save(output_path)


# 预设主题模板
THEME_PRESETS = {
    "business": {
        "name": "商务蓝",
        "color_scheme": {
            "primary_color": "#2E86AB",
            "secondary_color": "#1A365D",
            "accent_color": "#F18F01",
            "background_color": "#FFFFFF",
            "text_color": "#333333",
            "title_color": "#1A1A2E"
        }
    },
    "tech": {
        "name": "科技紫",
        "color_scheme": {
            "primary_color": "#667eea",
            "secondary_color": "#764ba2",
            "accent_color": "#00D9FF",
            "background_color": "#0F0F23",
            "text_color": "#E0E0E0",
            "title_color": "#FFFFFF"
        }
    },
    "nature": {
        "name": "自然绿",
        "color_scheme": {
            "primary_color": "#2D6A4F",
            "secondary_color": "#40916C",
            "accent_color": "#95D5B2",
            "background_color": "#FFFFFF",
            "text_color": "#333333",
            "title_color": "#1B4332"
        }
    },
    "warm": {
        "name": "温暖橙",
        "color_scheme": {
            "primary_color": "#E85D04",
            "secondary_color": "#DC2F02",
            "accent_color": "#FFBA08",
            "background_color": "#FFFBF5",
            "text_color": "#333333",
            "title_color": "#370617"
        }
    },
    "minimal": {
        "name": "极简黑白",
        "color_scheme": {
            "primary_color": "#333333",
            "secondary_color": "#666666",
            "accent_color": "#E63946",
            "background_color": "#FFFFFF",
            "text_color": "#333333",
            "title_color": "#000000"
        }
    },
    "ocean": {
        "name": "海洋蓝",
        "color_scheme": {
            "primary_color": "#0077B6",
            "secondary_color": "#023E8A",
            "accent_color": "#48CAE4",
            "background_color": "#CAF0F8",
            "text_color": "#03045E",
            "title_color": "#023E8A"
        }
    }
}


def apply_theme_preset(ppt_data: Dict[str, Any], theme_name: str) -> Dict[str, Any]:
    """
    应用预设主题到PPT数据
    
    Args:
        ppt_data: PPT数据字典
        theme_name: 主题名称
        
    Returns:
        更新后的PPT数据
    """
    if theme_name in THEME_PRESETS:
        ppt_data["color_scheme"] = THEME_PRESETS[theme_name]["color_scheme"]
    return ppt_data


if __name__ == "__main__":
    # 测试
    test_data = {
        "topic": "人工智能入门",
        "color_scheme": THEME_PRESETS["tech"]["color_scheme"],
        "slides": [
            {
                "slide_number": 1,
                "slide_type": "title",
                "title": "人工智能入门",
                "subtitle": "从基础到实践",
                "content": ["探索AI的无限可能"]
            },
            {
                "slide_number": 2,
                "slide_type": "content",
                "title": "什么是人工智能？",
                "content": [
                    "人工智能是计算机科学的一个分支",
                    "旨在创建能够执行需要人类智能的任务的系统",
                    "包括学习、推理、问题解决等能力"
                ]
            },
            {
                "slide_number": 3,
                "slide_type": "bullet_points",
                "title": "AI的核心技术",
                "content": [
                    "机器学习 - 让计算机从数据中学习",
                    "深度学习 - 使用神经网络进行复杂模式识别",
                    "自然语言处理 - 理解和生成人类语言",
                    "计算机视觉 - 让机器'看懂'图像和视频"
                ]
            },
            {
                "slide_number": 4,
                "slide_type": "two_column",
                "title": "AI的优势与挑战",
                "content": [
                    "高效处理大量数据",
                    "24/7不间断工作",
                    "快速决策能力",
                    "数据隐私问题",
                    "算法偏见风险",
                    "就业结构变化"
                ]
            },
            {
                "slide_number": 5,
                "slide_type": "summary",
                "title": "总结与展望",
                "content": [
                    "AI正在改变世界",
                    "持续学习是关键",
                    "拥抱技术变革"
                ]
            }
        ]
    }
    
    output_file = "test_ppt.pptx"
    create_ppt_from_data(test_data, output_file)
    print(f"PPT已保存到: {output_file}")
